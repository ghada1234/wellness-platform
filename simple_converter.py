#!/usr/bin/env python3
"""
Simple HTML to PPTX Converter
"""

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    print("✅ Required packages are available")
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Find Your Inner Peace"
    subtitle = slide.placeholders[1]
    subtitle.text = "Your Comprehensive Wellness & Meditation Platform"
    
    # Slide 2: Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Your Complete Wellness Companion"
    content = slide.placeholders[1]
    content.text = """• Mood Tracking & Analytics
• Guided Meditation & Breathing
• AI-Powered Insights
• Nutrition & Hydration Tracking
• Sleep & Activity Monitoring
• Journal & Reflection Tools"""
    
    # Slide 3: Mood Tracking
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Track Your Emotional Wellbeing"
    content = slide.placeholders[1]
    content.text = """• Daily mood logging with emojis
• Energy level tracking
• Influence factor identification
• Weekly mood analytics
• Trend visualization
• Personalized insights"""
    
    # Slide 4: Meditation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Find Peace Through Meditation"
    content = slide.placeholders[1]
    content.text = """• Guided breathing exercises
• Interactive meditation sessions
• Mindfulness practices
• Stress relief techniques
• Progress tracking
• Customizable sessions"""
    
    # Slide 5: AI Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "AI-Powered Wellness Insights"
    content = slide.placeholders[1]
    content.text = """• Personalized meal plans
• Nutrition analysis
• Wellness coaching
• Habit recommendations
• Progress insights
• Goal setting assistance"""
    
    # Save presentation
    prs.save('Find_Your_Inner_Peace_Presentation.pptx')
    print("✅ Presentation saved as: Find_Your_Inner_Peace_Presentation.pptx")
    print(f"📊 Total slides: {len(prs.slides)}")
    
except ImportError as e:
    print("❌ Missing packages. Installing...")
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    print("✅ Packages installed. Please run the script again.")
    
except Exception as e:
    print(f"❌ Error: {e}")

















